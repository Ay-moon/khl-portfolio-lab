from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "KHL_Architecture_Technique_Detaillee_40_Diapos.pptx"

COL = {
    "bg": "090F1B",
    "bg2": "121B2D",
    "panel": "0B172A",
    "panel2": "111F36",
    "text": "E8F2FF",
    "muted": "9EB4D2",
    "border": "97C4FF",
    "good": "4AD9A8",
    "setup": "235B7D",
    "platform": "2A4E86",
    "pipelines": "6F4E20",
    "intelligence": "5B3F85",
    "delivery": "2B5F5F",
    "global": "223F68",
}


def rgb(key: str) -> RGBColor:
    h = COL.get(key, key).lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def add_bg(slide) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb("bg")

    orb1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-0.95), Inches(-0.9), Inches(4.4), Inches(4.4))
    orb1.fill.solid()
    orb1.fill.fore_color.rgb = rgb("245EA4")
    orb1.fill.transparency = 0.64
    orb1.line.fill.background()

    orb2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.1), Inches(-1.15), Inches(4.25), Inches(4.25))
    orb2.fill.solid()
    orb2.fill.fore_color.rgb = rgb("238661")
    orb2.fill.transparency = 0.70
    orb2.line.fill.background()

    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.65), Inches(13.333), Inches(0.85))
    strip.fill.solid()
    strip.fill.fore_color.rgb = rgb("bg2")
    strip.fill.transparency = 0.35
    strip.line.fill.background()


def write_text(slide, x, y, w, h, text, size=12, color="text", bold=False, align=None, font="Manrope"):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = rgb(color)
    p.font.name = font
    if align is not None:
        p.alignment = align
    return tb


def header(slide, title: str, tag: str, idx: int, total: int) -> None:
    add_bg(slide)
    top = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.34), Inches(0.22), Inches(12.65), Inches(1.22))
    top.fill.solid()
    top.fill.fore_color.rgb = rgb("panel2")
    top.fill.transparency = 0.06
    top.line.color.rgb = rgb("border")
    top.line.width = Pt(1.0)

    write_text(slide, 0.56, 0.36, 9.6, 0.5, title, 22, "text", True, font="Archivo")
    write_text(slide, 0.56, 0.84, 9.6, 0.4, "Storytelling visuel du projet", 10, "muted")

    chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.45), Inches(0.47), Inches(2.2), Inches(0.46))
    chip.fill.solid()
    chip.fill.fore_color.rgb = rgb(tag)
    chip.fill.transparency = 0.05
    chip.line.color.rgb = rgb("border")
    chip.line.width = Pt(0.8)
    tf = chip.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"{tag.upper()}  {idx}/{total}"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = rgb("text")
    p.font.name = "Archivo"

    bar_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.55), Inches(12.2), Inches(0.08))
    bar_bg.fill.solid()
    bar_bg.fill.fore_color.rgb = rgb("panel2")
    bar_bg.line.fill.background()

    pct = max(0.03, min(1.0, idx / total))
    bar_fg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.55), Inches(12.2 * pct), Inches(0.08))
    bar_fg.fill.solid()
    bar_fg.fill.fore_color.rgb = rgb("good")
    bar_fg.line.fill.background()


def node(slide, x, y, w, h, label, color_key="panel2"):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = rgb(color_key)
    sh.fill.transparency = 0.04
    sh.line.color.rgb = rgb("border")
    sh.line.width = Pt(1.0)
    tf = sh.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = label
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = rgb("text")
    p.font.name = "Manrope"
    return sh


def arrow(slide, x, y, w=0.3, h=0.28):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    a.fill.solid()
    a.fill.fore_color.rgb = rgb("good")
    a.line.fill.background()


def slide_flow(prs: Presentation, idx: int, title: str, tag: str, items: list[str]) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, title, tag, idx, 40)

    panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(1.82), Inches(12.23), Inches(5.23))
    panel.fill.solid(); panel.fill.fore_color.rgb = rgb("panel"); panel.fill.transparency = 0.05
    panel.line.color.rgb = rgb("border"); panel.line.width = Pt(0.9)

    n = len(items)
    bw = 1.9 if n >= 6 else 2.2
    gap = 0.3
    total_w = n * bw + (n - 1) * gap
    sx = (13.333 - total_w) / 2
    y = 3.45

    for i, item in enumerate(items):
        x = sx + i * (bw + gap)
        node(s, x, y, bw, 1.1, item, "panel2")
        if i < n - 1:
            arrow(s, x + bw + 0.03, y + 0.41)


def slide_lanes(prs: Presentation, idx: int, title: str, tag: str, lanes: list[tuple[str, list[str]]]) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, title, tag, idx, 40)

    panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(1.82), Inches(12.23), Inches(5.23))
    panel.fill.solid(); panel.fill.fore_color.rgb = rgb("panel"); panel.fill.transparency = 0.05
    panel.line.color.rgb = rgb("border"); panel.line.width = Pt(0.9)

    lane_w = 12.0 / len(lanes)
    base_x = 0.66

    for i, (lt, nodes) in enumerate(lanes):
        lx = base_x + i * lane_w
        lane = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(lx), Inches(2.0), Inches(lane_w - 0.12), Inches(4.8))
        lane.fill.solid(); lane.fill.fore_color.rgb = rgb("panel2"); lane.fill.transparency = 0.1
        lane.line.color.rgb = rgb("border"); lane.line.width = Pt(0.7)
        write_text(s, lx + 0.08, 2.06, lane_w - 0.28, 0.3, lt, 10, "good", True, PP_ALIGN.CENTER)

        ny = 2.55
        for n in nodes[:4]:
            node(s, lx + 0.16, ny, lane_w - 0.44, 0.8, n, "panel")
            ny += 1.03


def slide_hub(prs: Presentation, idx: int, title: str, tag: str, center: str, orbit: list[str]) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, title, tag, idx, 40)

    panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(1.82), Inches(12.23), Inches(5.23))
    panel.fill.solid(); panel.fill.fore_color.rgb = rgb("panel"); panel.fill.transparency = 0.05
    panel.line.color.rgb = rgb("border"); panel.line.width = Pt(0.9)

    cx, cy = 6.67, 4.35
    node(s, cx - 1.2, cy - 0.55, 2.4, 1.1, center, "global")

    pos = [
        (2.0, 2.4), (4.2, 2.0), (9.1, 2.0), (11.2, 2.4),
        (2.0, 5.2), (4.2, 5.6), (9.1, 5.6), (11.2, 5.2),
    ]

    for i, label in enumerate(orbit[:8]):
        x, y = pos[i]
        node(s, x, y, 2.0, 0.8, label, "panel2")
        # connector line
        line = s.shapes.add_connector(1, Inches(cx), Inches(cy), Inches(x + 1.0), Inches(y + 0.4))
        line.line.color.rgb = rgb("border")
        line.line.width = Pt(0.8)


def slide_grid(prs: Presentation, idx: int, title: str, tag: str, cells: list[str]) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, title, tag, idx, 40)

    panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(1.82), Inches(12.23), Inches(5.23))
    panel.fill.solid(); panel.fill.fore_color.rgb = rgb("panel"); panel.fill.transparency = 0.05
    panel.line.color.rgb = rgb("border"); panel.line.width = Pt(0.9)

    cols, rows = 4, 3
    w, h = 2.75, 1.25
    gx, gy = 0.95, 2.25
    k = 0
    for r in range(rows):
        for cidx in range(cols):
            if k >= len(cells):
                break
            x = gx + cidx * 3.05
            y = gy + r * 1.5
            node(s, x, y, w, h, cells[k], "panel2")
            k += 1


def slide_timeline(prs: Presentation, idx: int, title: str, tag: str, steps: list[str]) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, title, tag, idx, 40)

    panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(1.82), Inches(12.23), Inches(5.23))
    panel.fill.solid(); panel.fill.fore_color.rgb = rgb("panel"); panel.fill.transparency = 0.05
    panel.line.color.rgb = rgb("border"); panel.line.width = Pt(0.9)

    y = 4.45
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(y), Inches(11.3), Inches(0.06))
    line.fill.solid(); line.fill.fore_color.rgb = rgb("border"); line.line.fill.background()

    n = len(steps)
    for i, st in enumerate(steps):
        x = 1.0 + (11.0 / max(1, n - 1)) * i
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x - 0.12), Inches(y - 0.13), Inches(0.24), Inches(0.24))
        dot.fill.solid(); dot.fill.fore_color.rgb = rgb("good"); dot.line.fill.background()
        node(s, x - 0.8, 3.05 if i % 2 == 0 else 5.0, 1.6, 0.75, st, "panel2")
        connector = s.shapes.add_connector(1, Inches(x), Inches(y), Inches(x), Inches(3.8 if i % 2 == 0 else 5.0))
        connector.line.color.rgb = rgb("border")
        connector.line.width = Pt(0.8)


def build_slides() -> list[dict]:
    return [
        {"kind": "cover", "title": "KHL Portfolio Lab", "tag": "global"},
        {"kind": "timeline", "title": "Acte 1 | Vision", "tag": "global", "steps": ["Probleme", "Vision", "Socle", "Flux", "Decision", "Impact"]},
        {"kind": "flow", "title": "Acte 1 | Story Arc", "tag": "global", "items": ["Source", "Modele", "Simulation", "Risque", "IA", "Delivery"]},
        {"kind": "hub", "title": "Acte 1 | Espace Projet", "tag": "global", "center": "KHL Lab", "orbit": ["Dev", "SQL STG", "SQL Gold", "Power BI", "CI", "Trading", "Analytics", "AI"]},
        {"kind": "lanes", "title": "Acte 1 | 5 Couches", "tag": "global", "lanes": [("Setup", [".env", ".venv", "Bootstrap"]), ("Platform", ["Dim", "Fact", "Views"]), ("Pipelines", ["Ingest", "Sim", "Calc"]), ("Intel", ["Risk", "AI", "Audit"]), ("Delivery", ["BI", "CI", "Release"])]},

        {"kind": "flow", "title": "Acte 2 | Setup", "tag": "setup", "items": ["Template", ".env", ".venv", "deps", "seed"]},
        {"kind": "grid", "title": "Acte 2 | Variables", "tag": "setup", "cells": ["SQL_SERVER", "SQL_DATABASE", "SQL_USER", "SQL_PASSWORD", "SQL_DRIVER", "SQL_AUTH_MODE", "SQL_DATABASE_STG", "SQL_SCHEMA_STG", "OPENAI_API_KEY", "PYTHON_EXE", "ENV", "DB_ENGINE"]},
        {"kind": "lanes", "title": "Acte 2 | Bootstrap Scripts", "tag": "setup", "lanes": [("Windows", ["bootstrap_local.ps1", "read .env", "pip install", "seed"]), ("Linux", ["bootstrap_local.sh", "read .env", "pip install", "seed"]), ("Result", ["DB ready", "Model ready", "Demo ready"]) ]},
        {"kind": "timeline", "title": "Acte 2 | Onboarding", "tag": "setup", "steps": ["Clone", "Bootstrap", "Seed", "Run", "Verify"]},

        {"kind": "flow", "title": "Acte 3 | Data Platform", "tag": "platform", "items": ["Dim", "Facts", "AI_*", "Views"]},
        {"kind": "grid", "title": "Acte 3 | Dimensions", "tag": "platform", "cells": ["DimDate", "DateKey", "Calendar", "MonthEnd", "DimSecurity", "Ticker", "AssetClass", "Currency", "DimPortfolio", "Code", "RiskProfile", "Inception"]},
        {"kind": "grid", "title": "Acte 3 | Facts", "tag": "platform", "cells": ["FactPrice", "FactTrades", "Positions", "PnL", "Risk", "DateKey", "SecurityKey", "PortfolioKey", "ClosePrice", "Quantity", "NAV", "VaR95"]},
        {"kind": "grid", "title": "Acte 3 | AI Tables", "tag": "platform", "cells": ["AI_DailyBriefing", "AI_Recommendations", "AI_WhatIf", "AI_AuditLog", "OutputJson", "Reasoning", "Scenario", "Status", "EventTs", "RequestId", "Detail", "ModelName"]},
        {"kind": "lanes", "title": "Acte 3 | Views", "tag": "platform", "lanes": [("Trades", ["vw_FactTradesEnriched", "Gross", "Net"]), ("Positions", ["vw_PositionSnapshot", "Weight", "Unrealized"]), ("Dashboard", ["vw_PortfolioDashboardDaily", "PnL", "Risk"]), ("AI", ["vw_AI_LatestRecommendations", "Latest", "Status"]) ]},
        {"kind": "hub", "title": "Acte 3 | Gold Core", "tag": "platform", "center": "SQL Gold", "orbit": ["DimDate", "DimSecurity", "DimPortfolio", "FactPrice", "FactTrades", "Positions", "PnL", "Risk"]},

        {"kind": "flow", "title": "Acte 4 | Ingestion Story", "tag": "pipelines", "items": ["STG", "Parse", "Canonical", "Dedup", "Merge"]},
        {"kind": "timeline", "title": "Acte 4 | Parsing", "tag": "pipelines", "steps": ["Raw", "Filter ACTION", "Parse FR", "Clean", "Ready"]},
        {"kind": "hub", "title": "Acte 4 | Canonical Security", "tag": "pipelines", "center": "Canonical Name", "orbit": ["libelle", "ss_jacent", "sous_jacent", "produit", "isin"]},
        {"kind": "flow", "title": "Acte 4 | Dedup Rule", "tag": "pipelines", "items": ["key(date,name)", "compare", "load_ts", "stg_id", "keep latest"]},
        {"kind": "lanes", "title": "Acte 4 | Merge FactPrice", "tag": "pipelines", "lanes": [("Source", ["#SrcFactPrice", "DateKey", "SecurityKey"]), ("Target", ["FactPrice", "match", "insert/update"]), ("Output", ["inserted", "updated", "source_system"]) ]},

        {"kind": "flow", "title": "Acte 5 | Trading Story", "tag": "pipelines", "items": ["Price Grid", "Strategy", "Orders", "Execute", "Snapshot", "SQL"]},
        {"kind": "grid", "title": "Acte 5 | Domain Objects", "tag": "pipelines", "cells": ["MarketOrder", "PositionState", "ExecutedTrade", "PositionSnapshot", "BUY/SELL", "Quantity", "AvgCost", "Fee", "Slippage", "CashFlow", "Weight", "PnL"]},
        {"kind": "timeline", "title": "Acte 5 | Pricing", "tag": "pipelines", "steps": ["Seed", "Business Days", "Drift", "Noise", "Price"]},
        {"kind": "flow", "title": "Acte 5 | Rotation Logic", "tag": "pipelines", "items": ["Pick Ticker", "BUY zone", "SELL check", "SPY pulse", "Orders"]},
        {"kind": "lanes", "title": "Acte 5 | run_mvp", "tag": "pipelines", "lanes": [("Prepare", ["ensure date", "ensure portfolio", "ensure securities"]), ("Run", ["execute trades", "build snapshots", "collect"]), ("Persist", ["purge old", "insert trades", "insert positions"]) ]},

        {"kind": "flow", "title": "Acte 6 | Performance Story", "tag": "intelligence", "items": ["Trades", "Cash", "Market Value", "NAV", "PnL"]},
        {"kind": "timeline", "title": "Acte 6 | Risk Story", "tag": "intelligence", "steps": ["Returns", "Vol20d", "Peak", "Drawdown", "VaR95", "Sharpe"]},
        {"kind": "lanes", "title": "Acte 6 | Guardrails", "tag": "intelligence", "lanes": [("Check", ["missing price?", "position > 0"]), ("Fail Fast", ["stop job", "show sample"]), ("Resume", ["ingest prices", "rerun calc"]) ]},
        {"kind": "flow", "title": "Acte 6 | Rerun", "tag": "intelligence", "items": ["select range", "purge", "recompute", "insert"]},
        {"kind": "hub", "title": "Acte 6 | Portfolio Dashboard", "tag": "intelligence", "center": "vw_PortfolioDashboardDaily", "orbit": ["NAV", "DailyPnL", "CumPnL", "Vol", "Drawdown", "VaR", "Sharpe"]},

        {"kind": "flow", "title": "Acte 7 | AI Runtime", "tag": "intelligence", "items": ["resolve", "retrieve", "generate", "validate", "write", "audit"]},
        {"kind": "grid", "title": "Acte 7 | Context Pack", "tag": "intelligence", "cells": ["headline", "positions", "latest recos", "date_key", "portfolio_code", "nav", "daily_pnl", "return_pct", "vol", "drawdown", "var95", "sharpe"]},
        {"kind": "lanes", "title": "Acte 7 | Commands", "tag": "intelligence", "lanes": [("daily-briefing", ["regime", "summary", "assumptions"]), ("recommendations", ["action", "confidence", "target_weight"]), ("what-if", ["scenario", "delta vol", "delta return"]) ]},
        {"kind": "flow", "title": "Acte 7 | Constraints", "tag": "intelligence", "items": ["schema", "actions", "limits", "ban phrases", "postcheck"]},
        {"kind": "hub", "title": "Acte 7 | AI Writeback", "tag": "intelligence", "center": "AI_AuditLog", "orbit": ["DailyBriefing", "Recommendations", "WhatIf", "status", "component", "detail"]},

        {"kind": "flow", "title": "Acte 8 | Delivery", "tag": "delivery", "items": ["Gold Views", "Power BI", "Insights", "Decision"]},
        {"kind": "grid", "title": "Acte 8 | BI Pages", "tag": "delivery", "cells": ["Overview", "Performance", "Risk", "AI Insights", "Trend", "Drilldown", "Exposure", "PnL", "VaR", "Drawdown", "Reco", "What-if"]},
        {"kind": "lanes", "title": "Acte 8 | Quality", "tag": "delivery", "lanes": [("Unit tests", ["simulator", "performance", "ingestion", "ai rules"]), ("CI Python", ["compileall", "unittest"]), ("CI SQL", ["validate_sql.py"]), ("Release", ["tag v*", "changelog", "zip"]) ]},
        {"kind": "timeline", "title": "Acte 8 | Roadmap", "tag": "delivery", "steps": ["Now", "OpenAI provider", "Prompt V2", "Docker", "Orchestration", "Prod"]},
        {"kind": "flow", "title": "Final | End-to-End", "tag": "delivery", "items": ["Setup", "Ingestion", "Trading", "Risk", "AI", "BI"]},
    ]


def build() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    specs = build_slides()
    assert len(specs) == 40

    for i, spec in enumerate(specs, start=1):
        kind = spec["kind"]

        if kind == "cover":
            s = prs.slides.add_slide(prs.slide_layouts[6])
            add_bg(s)
            hero = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(0.82), Inches(11.95), Inches(4.95))
            hero.fill.solid(); hero.fill.fore_color.rgb = rgb("panel2"); hero.fill.transparency = 0.04
            hero.line.color.rgb = rgb("border"); hero.line.width = Pt(1.3)
            write_text(s, 1.02, 1.23, 11.2, 1.45, "KHL Portfolio Lab\nStorytelling Architecture", 38, "text", True, font="Archivo")
            write_text(s, 1.02, 3.05, 10.6, 0.65, "40 slides | visual diagrams | low text", 16, "muted")
            chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.02), Inches(3.82), Inches(10.95), Inches(0.62))
            chip.fill.solid(); chip.fill.fore_color.rgb = rgb("global"); chip.fill.transparency = 0.14
            chip.line.color.rgb = rgb("border"); chip.line.width = Pt(0.8)
            tf = chip.text_frame; tf.clear(); p = tf.paragraphs[0]
            p.text = "Setup | Data Platform | Pipelines | Intelligence | Delivery"
            p.alignment = PP_ALIGN.CENTER; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = rgb("text"); p.font.name = "Archivo"
            continue

        if kind == "flow":
            slide_flow(prs, i, spec["title"], spec["tag"], spec["items"])
        elif kind == "lanes":
            slide_lanes(prs, i, spec["title"], spec["tag"], spec["lanes"])
        elif kind == "hub":
            slide_hub(prs, i, spec["title"], spec["tag"], spec["center"], spec["orbit"])
        elif kind == "grid":
            slide_grid(prs, i, spec["title"], spec["tag"], spec["cells"])
        elif kind == "timeline":
            slide_timeline(prs, i, spec["title"], spec["tag"], spec["steps"])
        else:
            raise RuntimeError(f"Unknown kind: {kind}")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Generated: {OUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
